#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_slide_with_title(prs, title, content_top=1.3):
    """Crear diapositiva con título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blanco

    # Fondo de título
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(52, 152, 219)
    title_shape.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide, content_top

def add_title_slide(prs):
    """Diapositiva de título principal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(44, 62, 80)

    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.text = "¿Niño o niña?\nLa percepción del género en las voces infantiles"
    for p in tf.paragraphs:
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_audio_slide(prs, slide_num, audios_info):
    """Diapositiva con audios"""
    title = f"Actividad inicial{''.join([' (II)', ' (III)'][slide_num-1:slide_num]) if slide_num > 1 else ': ¿Quién está hablando?'}"
    slide, y = add_slide_with_title(prs, title)

    if slide_num == 1:
        # Subtítulo
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
        tf = text_box.text_frame
        tf.text = "Escuchad estas voces y responded: ¿es un niño o una niña?"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += 0.8
    else:
        y += 0.3

    # Agregar audios
    for idx, (audio_num, audio_file) in enumerate(audios_info):
        col = idx % 2
        row = idx // 2

        x = 1.5 + (col * 5)
        audio_y = y + (row * 2)

        # Cuadro para el audio
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - 0.3), Inches(audio_y),
            Inches(3), Inches(1.5)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(236, 240, 241)
        rect.line.color.rgb = RGBColor(52, 152, 219)
        rect.line.width = Pt(2)

        # Texto "Audio X"
        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(audio_y + 0.1), Inches(2.6), Inches(0.4))
        tf = label_box.text_frame
        tf.text = f"Audio {audio_num}"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Insertar audio como objeto de medios
        if os.path.exists(audio_file):
            try:
                # PowerPoint embebe el audio y muestra un ícono
                slide.shapes.add_movie(
                    audio_file,
                    Inches(x + 0.5), Inches(audio_y + 0.5),
                    Inches(1), Inches(0.8),
                    poster_frame_image=None,
                    mime_type='audio/wav'
                )
            except Exception as e:
                print(f"Error con audio {audio_file}: {e}")
                # Fallback: indicador visual
                audio_text = slide.shapes.add_textbox(Inches(x), Inches(audio_y + 0.6), Inches(2.2), Inches(0.6))
                tf = audio_text.text_frame
                tf.text = f"🔊 {audio_file}"
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Pregunta final solo en última diapositiva de audios
    if slide_num == 3:
        question_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        tf = question_box.text_frame
        tf.text = "Pregunta: ¿habéis podido identificar el género de cada voz? ¿Con qué grado de certeza?"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(192, 57, 43)
        p.alignment = PP_ALIGN.CENTER

def add_text_slide(prs, title, subtitle, bullets):
    """Diapositiva con texto y viñetas"""
    slide, y = add_slide_with_title(prs, title)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
        tf = sub_box.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
        y += 0.7

    if bullets:
        for bullet in bullets:
            bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.5))
            tf = bullet_box.text_frame
            tf.text = bullet
            p = tf.paragraphs[0]
            p.font.size = Pt(22)
            p.level = 0
            y += 0.5

def add_table_slide(prs, title):
    """Diapositiva con tabla de datos"""
    slide, y = add_slide_with_title(prs, title)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    tf.text = "Parámetros de las grabaciones que escuchasteis"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
    y += 0.6

    # Crear tabla
    rows, cols = 7, 7
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(y), Inches(9.4), Inches(3.5)).table

    # Datos
    headers = ["Grabación", "Palabras", "Vocales", "Tono (Hz)", "F1 (Hz)", "F2 (Hz)", "F3 (Hz)"]
    data = [
        ["Niña 1", "4", "28", "300±38", "678±206", "1603±682", "2918±494"],
        ["Niña 2", "1", "16", "259±39", "702±186", "1376±339", "2568±583"],
        ["Niña 3", "3", "13", "238±23", "687±130", "1220±356", "2844±633"],
        ["Niño 1", "7", "23", "280±34", "660±132", "1741±494", "2697±406"],
        ["Niño 2", "2", "9", "268±39", "666±57", "1750±292", "2879±500"],
        ["Niño 3", "5", "15", "302±44", "681±140", "1598±501", "2800±573"],
    ]

    # Headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Datos
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            p.alignment = PP_ALIGN.CENTER

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

def add_image_slide(prs, title, subtitle, image_path, interpretation):
    """Diapositiva con imagen"""
    slide, y = add_slide_with_title(prs, title)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
    y += 0.6

    # Imagen
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(2), Inches(y), height=Inches(4.2))
        y += 4.4

    # Interpretación
    interp_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = interp_box.text_frame
    tf.text = f"Interpretación: {interpretation}"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(127, 140, 141)

def add_conclusion_slide(prs, title, subtitle, points):
    """Diapositiva de conclusión"""
    slide, y = add_slide_with_title(prs, title)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)
    y += 0.8

    for point in points:
        point_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.5))
        tf = point_box.text_frame
        tf.text = f"• {point}"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        y += 0.6

# Crear presentación
prs = create_presentation()

# Diapositiva 1: Título
add_title_slide(prs)

# Diapositivas 2-4: Audios
add_audio_slide(prs, 1, [(1, "audio_ninio_2.wav"), (2, "audio_ninia_1.wav")])
add_audio_slide(prs, 2, [(3, "audio_ninio_3.wav"), (4, "audio_ninia3.wav")])
add_audio_slide(prs, 3, [(5, "audio_ninio_1.wav"), (6, "audio_ninia_2.wav")])

# Diapositiva 5: El enigma
slide, y = add_slide_with_title(prs, "El enigma científico")
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = text_box.text_frame
tf.text = "Lo que dice la investigación"
tf.paragraphs[0].font.size = Pt(26)
tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
y += 0.7

text_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.8))
tf = text_box2.text_frame
tf.text = "Según Funk & Simpson (2023), identificamos el género de voces infantiles con una precisión del 70-84%, muy por encima del azar, pero:"
tf.paragraphs[0].font.size = Pt(22)
y += 1

quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(7.6), Inches(1.2))
tf = quote_box.text_frame
tf.text = '"Las diferencias en el aparato fonador entre niños y niñas antes de la pubertad son prácticamente inexistentes"\n\n— Fitch & Giedd (1999)'
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = RGBColor(52, 152, 219)
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y-0.1), Inches(8), Inches(1.4))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(236, 240, 241)
rect.line.color.rgb = RGBColor(52, 152, 219)
slide.shapes._spTree.remove(rect._element)
slide.shapes._spTree.insert(2, rect._element)
y += 1.6

question_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = question_box.text_frame
tf.text = "Entonces, ¿cómo lo hacemos?"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 57, 43)
p.alignment = PP_ALIGN.CENTER

# Diapositiva 6: Tabla de datos
add_table_slide(prs, "Los datos acústicos")

# Diapositiva 7: Observación clave
slide, y = add_slide_with_title(prs, "Los datos acústicos (II)")
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = sub_box.text_frame
tf.text = "Observación clave"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)

obs_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
tf = obs_box.text_frame
tf.text = "Los rangos se solapan completamente.\n\nNo hay diferencias estadísticamente significativas."
for p in tf.paragraphs:
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.alignment = PP_ALIGN.CENTER

# Diapositiva 8: El tono
add_text_slide(prs, "Entendiendo la acústica de la voz", "El tono (frecuencia fundamental, F₀)", [
    "la \"altura\" de la voz (grave o aguda)",
    "producido por la vibración de las cuerdas vocales",
    "en niños prepuberales: 200-350 Hz (similar en ambos géneros)",
    "para comparar: adultos varones ~120 Hz, mujeres adultas ~220 Hz"
])

# Diapositiva 9: Formantes
slide, y = add_slide_with_title(prs, "Entendiendo la acústica de la voz (II)")
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = sub_box.text_frame
tf.text = "Los formantes (F1, F2, F3)"
tf.paragraphs[0].font.size = Pt(26)
tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
y += 0.7

bullets = [
    "frecuencias de resonancia del tracto vocal",
    "determinan la calidad de las vocales (/a/, /e/, /i/, /o/, /u/)",
    "relacionados con la longitud del tracto vocal"
]
for bullet in bullets:
    b_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.45))
    tf = b_box.text_frame
    tf.text = f"• {bullet}"
    tf.paragraphs[0].font.size = Pt(22)
    y += 0.5

y += 0.3
details = [
    "F1: apertura de la boca (bajo = cerrada /i/, alto = abierta /a/)",
    "F2: posición de la lengua (bajo = posterior /u/, alto = anterior /i/)",
    "F3: configuración más compleja del tracto vocal"
]
for detail in details:
    d_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.45))
    tf = d_box.text_frame
    tf.text = detail
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    y += 0.5

# Diapositivas 10-11: Imágenes
add_image_slide(prs, "Las visualizaciones acústicas", "Espacio vocálico F1-F2",
                "vowel_spaces_overlap_small.png",
                "las elipses muestran la distribución de las vocales de cada hablante. El solapamiento es evidente.")

add_image_slide(prs, "Las visualizaciones acústicas (II)", "Distribución del tono",
                "gender_comparison_statistical_small.png",
                "las barras de error muestran que los rangos de tono son muy similares entre niños y niñas.")

# Diapositiva 12: Paradoja - Percepción
slide, y = add_slide_with_title(prs, "La paradoja: ¿cómo diferenciamos entonces?")
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
tf = sub_box.text_frame
tf.text = "Lo que sabemos de la percepción"
tf.paragraphs[0].font.size = Pt(26)
tf.paragraphs[0].font.color.rgb = RGBColor(85, 85, 85)
y += 0.6

# Barreda & Assmann
ref_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.4))
tf = ref_box.text_frame
tf.text = "Barreda & Assmann (2021)"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(41, 128, 185)
y += 0.5

quote_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.9))
tf = quote_box.text_frame
tf.text = '"La percepción del género y la edad del hablante están entrelazadas. Los oyentes usan información sobre la edad para informar sus juicios de género"'
p = tf.paragraphs[0]
p.font.size = Pt(19)
p.font.italic = True
p.font.color.rgb = RGBColor(52, 152, 219)
y += 1

impl_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.5), Inches(0.5))
tf = impl_box.text_frame
tf.text = "Implicación: el contexto y las expectativas importan."
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True

# Diapositiva 13: Funk & Simpson
add_text_slide(prs, "Lo que sabemos de la percepción (II)", "Funk & Simpson (2023) - Identificaron varios factores clave:", [
    "Pitch como predictor principal (aunque con mucho solapamiento)",
    "Espectro de sibilantes (/s/, /z/): los niños tienden a producirlas con energía más baja",
    "Correlación con conformidad de género: los niños que expresan mayor conformidad con estereotipos de género muestran diferencias más marcadas"
])

# Diapositivas 14-16: Factores
add_text_slide(prs, "La respuesta: no es solo la anatomía", "Factor 1: diferencias comportamentales", [
    "Desde los 2-3 años, los niños internalizan estereotipos de género",
    "Pueden modificar voluntariamente su voz para sonar más \"masculinos\" o \"femeninos\"",
    "Cartei et al. (2019): niños de 6-10 años pueden controlar la expresión de masculinidad/feminidad en su voz"
])

add_text_slide(prs, "La respuesta: no es solo la anatomía (II)", "Factor 2: información prosódica", [
    "Patrones de entonación",
    "Ritmo del habla",
    "Variabilidad temporal y espectral",
    "Mucho más evidente en frases completas que en sílabas aisladas"
])

add_text_slide(prs, "La respuesta: no es solo la anatomía (III)", "Factor 3: información contextual", [
    "Duración del estímulo (mejor en oraciones que en vocales aisladas)",
    "Conocimiento de la edad aproximada del hablante",
    "Expectativas culturales"
])

# Diapositivas 17-19: Conclusiones
add_conclusion_slide(prs, "Conclusiones", "Las diferencias acústicas prepuberales son sutiles", [
    "No hay dimorfismo sexual anatómico significativo antes de la pubertad",
    "Los parámetros acústicos básicos (tono, formantes) se solapan completamente"
])

add_conclusion_slide(prs, "Conclusiones (II)", "Pero la percepción es robusta", [
    "Identificamos correctamente el género en ~70-80% de los casos",
    "La precisión mejora con más contexto (oraciones vs sílabas aisladas)"
])

add_conclusion_slide(prs, "Conclusiones (III)", "La voz como práctica social", [
    "Los niños aprenden y practican patrones de habla asociados a su género",
    "La voz no solo refleja anatomía, sino identidad de género",
    "Implicaciones: desarrollo del lenguaje, identidad de género en la infancia, terapia de voz"
])

# Diapositiva 20: Reflexión final
slide, y = add_slide_with_title(prs, "Reflexión final")
question_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = question_box.text_frame
tf.text = 'La pregunta no es solo "¿cómo diferenciamos?"'
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
y += 0.8

answer_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
tf = answer_box.text_frame
tf.text = "Es también: ¿Qué nos dice esto sobre cómo se construye el género?"
p = tf.paragraphs[0]
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 57, 43)
p.alignment = PP_ALIGN.CENTER
y += 1

concepts = [
    "Es performativo: se practica y se expresa",
    "Es perceptivo: lo interpretamos con expectativas culturales",
    "Es dinámico: evoluciona con el desarrollo"
]
for concept in concepts:
    c_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
    tf = c_box.text_frame
    tf.text = f"• {concept}"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    y += 0.6

# Diapositiva 21: Preguntas para debate
slide, y = add_slide_with_title(prs, "Preguntas para el debate")
questions = [
    "¿Creéis que los niños son conscientes de que modifican su voz para sonar más \"masculinos\" o \"femeninos\"?",
    "Si las diferencias anatómicas son mínimas, ¿de dónde aprenden los niños estos patrones vocales?",
    "¿Qué implicaciones tiene esto para nuestra comprensión del género como constructo social vs biológico?",
    "¿Debería esto cambiar nuestra aproximación a la terapia de voz para niños transgénero?"
]

for idx, question in enumerate(questions):
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.7))
    tf = q_box.text_frame
    tf.text = f"{idx+1}. {question}"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    y += 0.8

# Guardar
prs.save('presentacion.pptx')
print(f"Presentación generada: presentacion.pptx")
print(f"Total de diapositivas: {len(prs.slides)}")
