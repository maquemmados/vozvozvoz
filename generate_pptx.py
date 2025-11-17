#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

RED = RGBColor(192, 57, 43)
DARK = RGBColor(52, 73, 94)
LIGHT = RGBColor(245, 245, 245)

def add_header(slide, title):
    """Barra roja superior"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.85))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RED
    rect.line.fill.background()

    txt = slide.shapes.add_textbox(Inches(0.2), Inches(0.18), Inches(9.6), Inches(0.5))
    txt.text_frame.text = title
    txt.text_frame.paragraphs[0].font.size = Pt(30)
    txt.text_frame.paragraphs[0].font.bold = True
    txt.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# PORTADA
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(3))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "¿Niño o niña?"
p1.font.size = Pt(52)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "La percepción del género en las voces infantiles"
p2.font.size = Pt(36)
p2.font.color.rgb = RED
p2.alignment = PP_ALIGN.CENTER

# AUDIOS 1-2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Actividad inicial: ¿Quién está hablando?")

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.6))
sub.text_frame.word_wrap = True
sub.text_frame.text = "Escuchad estas voces y responded: ¿es un niño o una niña?"
sub.text_frame.paragraphs[0].font.size = Pt(24)
sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
sub.text_frame.paragraphs[0].font.bold = True

# Audio 1
b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(2), Inches(3.3), Inches(4.5))
b1.fill.solid()
b1.fill.fore_color.rgb = LIGHT
b1.line.color.rgb = RED
b1.line.width = Pt(4)

l1 = slide.shapes.add_textbox(Inches(1.6), Inches(2.3), Inches(2.9), Inches(0.6))
l1.text_frame.text = "Audio 1"
l1.text_frame.paragraphs[0].font.size = Pt(30)
l1.text_frame.paragraphs[0].font.bold = True
l1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninio_2.wav"):
    slide.shapes.add_movie("audio_ninio_2.wav", Inches(2.15), Inches(3.2), Inches(1.8), Inches(1.8))

# Audio 2
b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2), Inches(3.3), Inches(4.5))
b2.fill.solid()
b2.fill.fore_color.rgb = LIGHT
b2.line.color.rgb = RED
b2.line.width = Pt(4)

l2 = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(2.9), Inches(0.6))
l2.text_frame.text = "Audio 2"
l2.text_frame.paragraphs[0].font.size = Pt(30)
l2.text_frame.paragraphs[0].font.bold = True
l2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninia_1.wav"):
    slide.shapes.add_movie("audio_ninia_1.wav", Inches(6.05), Inches(3.2), Inches(1.8), Inches(1.8))

# AUDIOS 3-4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Actividad inicial (II)")

b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.5), Inches(3.3), Inches(4.8))
b3.fill.solid()
b3.fill.fore_color.rgb = LIGHT
b3.line.color.rgb = RED
b3.line.width = Pt(4)

l3 = slide.shapes.add_textbox(Inches(1.6), Inches(1.8), Inches(2.9), Inches(0.6))
l3.text_frame.text = "Audio 3"
l3.text_frame.paragraphs[0].font.size = Pt(30)
l3.text_frame.paragraphs[0].font.bold = True
l3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninio_3.wav"):
    slide.shapes.add_movie("audio_ninio_3.wav", Inches(2.15), Inches(2.8), Inches(1.8), Inches(1.8))

b4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.5), Inches(3.3), Inches(4.8))
b4.fill.solid()
b4.fill.fore_color.rgb = LIGHT
b4.line.color.rgb = RED
b4.line.width = Pt(4)

l4 = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(2.9), Inches(0.6))
l4.text_frame.text = "Audio 4"
l4.text_frame.paragraphs[0].font.size = Pt(30)
l4.text_frame.paragraphs[0].font.bold = True
l4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninia3.wav"):
    slide.shapes.add_movie("audio_ninia3.wav", Inches(6.05), Inches(2.8), Inches(1.8), Inches(1.8))

# AUDIOS 5-6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Actividad inicial (III)")

b5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.3), Inches(3.3), Inches(4))
b5.fill.solid()
b5.fill.fore_color.rgb = LIGHT
b5.line.color.rgb = RED
b5.line.width = Pt(4)

l5 = slide.shapes.add_textbox(Inches(1.6), Inches(1.6), Inches(2.9), Inches(0.6))
l5.text_frame.text = "Audio 5"
l5.text_frame.paragraphs[0].font.size = Pt(30)
l5.text_frame.paragraphs[0].font.bold = True
l5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninio_1.wav"):
    slide.shapes.add_movie("audio_ninio_1.wav", Inches(2.15), Inches(2.5), Inches(1.8), Inches(1.8))

b6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.3), Inches(3.3), Inches(4))
b6.fill.solid()
b6.fill.fore_color.rgb = LIGHT
b6.line.color.rgb = RED
b6.line.width = Pt(4)

l6 = slide.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(2.9), Inches(0.6))
l6.text_frame.text = "Audio 6"
l6.text_frame.paragraphs[0].font.size = Pt(30)
l6.text_frame.paragraphs[0].font.bold = True
l6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if os.path.exists("audio_ninia_2.wav"):
    slide.shapes.add_movie("audio_ninia_2.wav", Inches(6.05), Inches(2.5), Inches(1.8), Inches(1.8))

# Pregunta
q = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.3))
q.text_frame.word_wrap = True
q.text_frame.text = "Pregunta: ¿habéis podido identificar el género de cada voz?\n¿Con qué grado de certeza?"
q.text_frame.paragraphs[0].font.size = Pt(24)
q.text_frame.paragraphs[0].font.bold = True
q.text_frame.paragraphs[0].font.color.rgb = RED
q.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ENIGMA
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "El enigma científico")

y = 1.05

sub = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
sub.text_frame.text = "Lo que dice la investigación"
sub.text_frame.paragraphs[0].font.size = Pt(26)
sub.text_frame.paragraphs[0].font.color.rgb = DARK
sub.text_frame.paragraphs[0].font.bold = True
y += 0.6

txt = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(1))
txt.text_frame.word_wrap = True
txt.text_frame.text = "Según Funk & Simpson (2023), identificamos el género de voces infantiles con una precisión del 70-84%, muy por encima del azar, pero:"
txt.text_frame.paragraphs[0].font.size = Pt(23)
y += 1.1

qbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(8.4), Inches(1.7))
qbox.fill.solid()
qbox.fill.fore_color.rgb = LIGHT
qbox.line.color.rgb = RED
qbox.line.width = Pt(2)

qtxt = slide.shapes.add_textbox(Inches(1.2), Inches(y+0.15), Inches(7.6), Inches(1.4))
qtxt.text_frame.word_wrap = True
qtxt.text_frame.text = '"Las diferencias en el aparato fonador entre niños y niñas antes de la pubertad son prácticamente inexistentes"\n\n— Fitch & Giedd (1999)'
qtxt.text_frame.paragraphs[0].font.size = Pt(21)
qtxt.text_frame.paragraphs[0].font.italic = True
qtxt.text_frame.paragraphs[0].font.color.rgb = DARK
y += 1.9

final = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.9))
final.text_frame.word_wrap = True
final.text_frame.text = "Entonces, ¿cómo lo hacemos?"
final.text_frame.paragraphs[0].font.size = Pt(32)
final.text_frame.paragraphs[0].font.bold = True
final.text_frame.paragraphs[0].font.color.rgb = RED
final.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# TABLA
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Los datos acústicos")

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.4))
sub.text_frame.text = "Parámetros de las grabaciones que escuchasteis"
sub.text_frame.paragraphs[0].font.size = Pt(22)
sub.text_frame.paragraphs[0].font.color.rgb = DARK

table = slide.shapes.add_table(7, 7, Inches(0.15), Inches(1.5), Inches(9.7), Inches(4.8)).table

headers = ["Grabación", "Palabras", "Vocales", "Tono (Hz)", "F1 (Hz)", "F2 (Hz)", "F3 (Hz)"]
data = [
    ["Niña 1", "4", "28", "300±38", "678±206", "1603±682", "2918±494"],
    ["Niña 2", "1", "16", "259±39", "702±186", "1376±339", "2568±583"],
    ["Niña 3", "3", "13", "238±23", "687±130", "1220±356", "2844±633"],
    ["Niño 1", "7", "23", "280±34", "660±132", "1741±494", "2697±406"],
    ["Niño 2", "2", "9", "268±39", "666±57", "1750±292", "2879±500"],
    ["Niño 3", "5", "15", "302±44", "681±140", "1598±501", "2800±573"],
]

for col_idx, h in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = RED
    cell.text_frame.paragraphs[0].font.size = Pt(17)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

for row_idx, row in enumerate(data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(252, 252, 252)

obs = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
obs.text_frame.word_wrap = True
obs.text_frame.text = "Los rangos se solapan completamente"
obs.text_frame.paragraphs[0].font.size = Pt(26)
obs.text_frame.paragraphs[0].font.bold = True
obs.text_frame.paragraphs[0].font.color.rgb = RED
obs.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Continúa con el resto de diapositivas de manera similar...
# Por brevedad, incluyo solo funciones para las restantes

def text_slide(title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)

    y = 1.05
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.65))
        sub.text_frame.word_wrap = True
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(25)
        sub.text_frame.paragraphs[0].font.color.rgb = DARK
        sub.text_frame.paragraphs[0].font.bold = True
        y += 0.75

    for bullet in bullets:
        b = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.7))
        b.text_frame.word_wrap = True
        b.text_frame.text = f"• {bullet}"
        b.text_frame.paragraphs[0].font.size = Pt(23)
        y += 0.7

def image_slide(title, subtitle, img, interp):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title)

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.5))
    sub.text_frame.text = subtitle
    sub.text_frame.paragraphs[0].font.size = Pt(23)
    sub.text_frame.paragraphs[0].font.color.rgb = DARK

    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(1.5), Inches(1.7), height=Inches(4.6))

    interp_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    interp_box.text_frame.word_wrap = True
    interp_box.text_frame.text = f"Interpretación: {interp}"
    interp_box.text_frame.paragraphs[0].font.size = Pt(19)
    interp_box.text_frame.paragraphs[0].font.italic = True
    interp_box.text_frame.paragraphs[0].font.color.rgb = DARK

text_slide("Entendiendo la acústica de la voz", "El tono (frecuencia fundamental, F₀)", [
    'la "altura" de la voz (grave o aguda)',
    "producido por la vibración de las cuerdas vocales",
    "en niños prepuberales: 200-350 Hz (similar en ambos géneros)",
    "para comparar: adultos varones ~120 Hz, mujeres adultas ~220 Hz"
])

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Entendiendo la acústica de la voz (II)")
sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.6))
sub.text_frame.text = "Los formantes (F1, F2, F3)"
sub.text_frame.paragraphs[0].font.size = Pt(25)
sub.text_frame.paragraphs[0].font.color.rgb = DARK
sub.text_frame.paragraphs[0].font.bold = True

y = 1.75
for txt in ["frecuencias de resonancia del tracto vocal",
            "determinan la calidad de las vocales (/a/, /e/, /i/, /o/, /u/)",
            "relacionados con la longitud del tracto vocal"]:
    b = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    b.text_frame.word_wrap = True
    b.text_frame.text = f"• {txt}"
    b.text_frame.paragraphs[0].font.size = Pt(23)
    y += 0.65

y += 0.2
for det in ["F1: apertura de la boca (bajo = cerrada /i/, alto = abierta /a/)",
            "F2: posición de la lengua (bajo = posterior /u/, alto = anterior /i/)",
            "F3: configuración más compleja del tracto vocal"]:
    d = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
    d.text_frame.word_wrap = True
    d.text_frame.text = det
    d.text_frame.paragraphs[0].font.size = Pt(22)
    d.text_frame.paragraphs[0].font.bold = True
    d.text_frame.paragraphs[0].font.color.rgb = RED
    y += 0.65

image_slide("Las visualizaciones acústicas", "Espacio vocálico F1-F2",
            "vowel_spaces_overlap_small.png",
            "las elipses muestran la distribución de las vocales de cada hablante. El solapamiento es evidente.")

image_slide("Las visualizaciones acústicas (II)", "Distribución del tono",
            "gender_comparison_statistical_small.png",
            "las barras de error muestran que los rangos de tono son muy similares entre niños y niñas.")

# Percepción
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "La paradoja: ¿cómo diferenciamos entonces?")
sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.5))
sub.text_frame.text = "Lo que sabemos de la percepción - Barreda & Assmann (2021)"
sub.text_frame.paragraphs[0].font.size = Pt(23)
sub.text_frame.paragraphs[0].font.color.rgb = DARK
sub.text_frame.paragraphs[0].font.bold = True

qbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.6))
qbox.fill.solid()
qbox.fill.fore_color.rgb = LIGHT
qbox.line.color.rgb = RED
qbox.line.width = Pt(2)

qtxt = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(7.6), Inches(1.2))
qtxt.text_frame.word_wrap = True
qtxt.text_frame.text = '"La percepción del género y la edad del hablante están entrelazadas. Los oyentes usan información sobre la edad para informar sus juicios de género"'
qtxt.text_frame.paragraphs[0].font.size = Pt(21)
qtxt.text_frame.paragraphs[0].font.italic = True

impl = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(8.4), Inches(0.7))
impl.text_frame.word_wrap = True
impl.text_frame.text = "Implicación: el contexto y las expectativas importan."
impl.text_frame.paragraphs[0].font.size = Pt(26)
impl.text_frame.paragraphs[0].font.bold = True
impl.text_frame.paragraphs[0].font.color.rgb = RED

text_slide("Lo que sabemos de la percepción (II)", "Funk & Simpson (2023)", [
    "Pitch como predictor principal (aunque con mucho solapamiento)",
    "Espectro de sibilantes (/s/, /z/): los niños tienden a producirlas con energía más baja",
    "Correlación con conformidad de género: los niños que expresan mayor conformidad muestran diferencias más marcadas"
])

text_slide("La respuesta: no es solo la anatomía", "Factor 1: diferencias comportamentales", [
    "Desde los 2-3 años, los niños internalizan estereotipos de género",
    'Pueden modificar voluntariamente su voz para sonar más "masculinos" o "femeninos"',
    "Cartei et al. (2019): niños de 6-10 años pueden controlar la expresión de masculinidad/feminidad"
])

text_slide("La respuesta: no es solo la anatomía (II)", "Factor 2: información prosódica", [
    "Patrones de entonación",
    "Ritmo del habla",
    "Variabilidad temporal y espectral",
    "Mucho más evidente en frases completas que en sílabas aisladas"
])

text_slide("La respuesta: no es solo la anatomía (III)", "Factor 3: información contextual", [
    "Duración del estímulo (mejor en oraciones que en vocales aisladas)",
    "Conocimiento de la edad aproximada del hablante",
    "Expectativas culturales"
])

text_slide("Conclusiones", "Las diferencias acústicas prepuberales son sutiles", [
    "No hay dimorfismo sexual anatómico significativo antes de la pubertad",
    "Los parámetros acústicos básicos (tono, formantes) se solapan completamente"
])

text_slide("Conclusiones (II)", "Pero la percepción es robusta", [
    "Identificamos correctamente el género en ~70-80% de los casos",
    "La precisión mejora con más contexto (oraciones vs sílabas aisladas)"
])

text_slide("Conclusiones (III)", "La voz como práctica social", [
    "Los niños aprenden y practican patrones de habla asociados a su género",
    "La voz no solo refleja anatomía, sino identidad de género",
    "Implicaciones: desarrollo del lenguaje, identidad de género, terapia de voz"
])

# Reflexión final
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Reflexión final")

q1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.7))
q1.text_frame.word_wrap = True
q1.text_frame.text = 'La pregunta no es solo "¿cómo diferenciamos?"'
q1.text_frame.paragraphs[0].font.size = Pt(29)
q1.text_frame.paragraphs[0].font.bold = True
q1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

q2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.9))
q2.text_frame.word_wrap = True
q2.text_frame.text = "Es también: ¿Qué nos dice esto sobre cómo se construye el género?"
q2.text_frame.paragraphs[0].font.size = Pt(27)
q2.text_frame.paragraphs[0].font.bold = True
q2.text_frame.paragraphs[0].font.color.rgb = RED
q2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

y = 3.2
for concept in ["Es performativo: se practica y se expresa",
                "Es perceptivo: lo interpretamos con expectativas culturales",
                "Es dinámico: evoluciona con el desarrollo"]:
    c = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.7))
    c.text_frame.word_wrap = True
    c.text_frame.text = f"• {concept}"
    c.text_frame.paragraphs[0].font.size = Pt(25)
    c.text_frame.paragraphs[0].font.bold = True
    y += 0.8

# Preguntas debate
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "Preguntas para el debate")

questions = [
    '¿Creéis que los niños son conscientes de que modifican su voz para sonar más "masculinos" o "femeninos"?',
    "Si las diferencias anatómicas son mínimas, ¿de dónde aprenden los niños estos patrones vocales?",
    "¿Qué implicaciones tiene esto para nuestra comprensión del género como constructo social vs biológico?",
    "¿Debería esto cambiar nuestra aproximación a la terapia de voz para niños transgénero?"
]

y = 1.3
for idx, q in enumerate(questions):
    qbox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    qbox.text_frame.word_wrap = True
    qbox.text_frame.text = f"{idx+1}. {q}"
    qbox.text_frame.paragraphs[0].font.size = Pt(21)
    y += 1.1

prs.save('presentacion.pptx')
print(f"✓ Presentación final: {len(prs.slides)} diapositivas")
